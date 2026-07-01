#!/usr/bin/env python3
"""
Presentation Generator — Bluestock MF Capstone
Uses python-pptx to generate the final 12-slide presentation.

Usage: python scripts/generate_presentation.py
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Paths
ROOT = Path(__file__).resolve().parent.parent
REPORTS = ROOT / 'reports'
CHARTS = REPORTS / 'charts'
REPORTS.mkdir(parents=True, exist_ok=True)

# Colors
DARK_BG = RGBColor(0x1a, 0x1a, 0x2e)
CARD_BG = RGBColor(0x16, 0x21, 0x3e)
RED_ACCENT = RGBColor(0xe9, 0x45, 0x60)
BLUE_ACCENT = RGBColor(0x0f, 0x34, 0x60)
WHITE = RGBColor(0xff, 0xff, 0xff)
GREY = RGBColor(0xa8, 0xa8, 0xb3)

def set_slide_background(slide, color):
    """Set the slide background color to solid RGB color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_header(slide, title_text):
    """Add standard section header to slide."""
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.33), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text.upper()
    p.font.name = 'Helvetica'
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RED_ACCENT
    
    # Bottom dividing line
    # (We can use a thin line or textbox with border, but keep it simple with text layout or formatting)

def add_textbox(slide, text, left, top, width, height, font_size=14, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    """Helper to add standard text block to slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    
    p = tf.paragraphs[0]
    p.alignment = align
    p.text = text
    p.font.name = 'Helvetica'
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    return txBox

def add_bullet_list(slide, items, left, top, width, height, font_size=12, color=WHITE):
    """Helper to add a bulleted text box to slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.name = 'Helvetica'
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.level = 0
        p.space_after = Pt(6)
    return txBox

def add_kpi_card(slide, label, value, left, top, width, height):
    """Draw a styled KPI box on the slide."""
    card = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE = 1
        left, top, width, height
    )
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = RED_ACCENT
    card.line.width = Pt(1.5)
    
    # Text
    tf = card.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = label.upper()
    p.font.name = 'Helvetica'
    p.font.size = Pt(10)
    p.font.color.rgb = GREY
    p.font.bold = True
    
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.text = value
    p2.font.name = 'Helvetica'
    p2.font.size = Pt(20)
    p2.font.color.rgb = WHITE
    p2.font.bold = True

def add_table(slide, rows_data, left, top, width, height, col_widths=None):
    """Draw a beautiful data table on slide."""
    rows = len(rows_data)
    cols = len(rows_data[0])
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    
    # Column Widths
    if col_widths:
        for idx, w in enumerate(col_widths):
            table.columns[idx].width = w
            
    # Style Cells
    for row_idx, row in enumerate(rows_data):
        for col_idx, text in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(text)
            
            # Text layout
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            p.font.name = 'Helvetica'
            p.font.size = Pt(10)
            
            if row_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = BLUE_ACCENT
                p.font.bold = True
                p.font.color.rgb = WHITE
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = CARD_BG if row_idx % 2 == 0 else DARK_BG
                p.font.color.rgb = WHITE
                
def main():
    print("Generating Bluestock_MF_Presentation.pptx...")
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # -------------------------------------------------------------------------
    # SLIDE 1: TITLE SLIDE
    # -------------------------------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_BG)
    
    # Top Accent Block
    top_line = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(0.2))
    top_line.fill.solid()
    top_line.fill.fore_color.rgb = RED_ACCENT
    top_line.line.fill.background()
    
    add_textbox(slide, "BLUESTOCK FINTECH", Inches(1.0), Inches(1.8), Inches(11.33), Inches(0.5), font_size=16, bold=True, color=RED_ACCENT)
    add_textbox(slide, "Mutual Fund Analytics Platform", Inches(1.0), Inches(2.4), Inches(11.33), Inches(1.5), font_size=38, bold=True, color=WHITE)
    add_textbox(slide, "End-to-End Data Engineering, ETL Pipeline & Interactive Dashboard", Inches(1.0), Inches(3.9), Inches(11.33), Inches(0.8), font_size=16, color=GREY)
    
    add_textbox(slide, "Prepared by: Adib Shaikh\nRole: Data Analyst Intern\nDate: June 2026", Inches(1.0), Inches(5.2), Inches(5.0), Inches(1.5), font_size=13, color=WHITE)
    
    # -------------------------------------------------------------------------
    # SLIDE 2: PROBLEM & OBJECTIVE
    # -------------------------------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_BG)
    add_header(slide, "The Problem & Objectives")
    
    probs = [
        "P1. Data Fragmentation: Mutual fund metrics split across disparate CSVs and APIs.",
        "P2. Return-Only Evaluation: Core fund analysis neglects costs and risk limits.",
        "P3. Lack of Benchmark Overlay: Active manager outperformance is rarely quantified.",
        "P4. Investor Retention Blind Spot: Irregular SIP payers are not tracked or flagged.",
        "P5. Reporting Delays: Processing and dashboard compilation take days, not seconds."
    ]
    add_textbox(slide, "THE CORE PROBLEMS", Inches(0.5), Inches(1.3), Inches(5.8), Inches(0.4), font_size=14, bold=True, color=RED_ACCENT)
    add_bullet_list(slide, probs, Inches(0.5), Inches(1.8), Inches(5.8), Inches(4.8), font_size=11)
    
    objs = [
        "1. Build programmatic ETL pipeline for 10 datasets.",
        "2. Normalise and store data in SQLite Star Schema.",
        "3. Optimize database query times with 6 index keys.",
        "4. Calculate 3yr CAGR, Sharpe, and Sortino risk ratios.",
        "5. Quantify alpha/beta and tracking error vs benchmarks.",
        "6. Model cohort size and SIP transaction gaps.",
        "7. Integrate a risk recommender in Streamlit dashboard.",
        "8. Export findings into visual final presentation decks."
    ]
    add_textbox(slide, "PROJECT OBJECTIVES", Inches(6.8), Inches(1.3), Inches(6.0), Inches(0.4), font_size=14, bold=True, color=RED_ACCENT)
    add_bullet_list(slide, objs, Inches(6.8), Inches(1.8), Inches(6.0), Inches(4.8), font_size=11)
    
    # -------------------------------------------------------------------------
    # SLIDE 3: DATA scale & SOURCES
    # -------------------------------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_BG)
    add_header(slide, "10 Real-World Datasets · 87,543 Rows")
    
    add_kpi_card(slide, "Total Funds", "40 Schemes", Inches(0.5), Inches(1.4), Inches(2.8), Inches(1.2))
    add_kpi_card(slide, "Time Scale", "4.5 Years", Inches(3.7), Inches(1.4), Inches(2.8), Inches(1.2))
    add_kpi_card(slide, "Transactions", "32,778 Rows", Inches(6.9), Inches(1.4), Inches(2.8), Inches(1.2))
    add_kpi_card(slide, "Total Database", "87,543 Rows", Inches(10.1), Inches(1.4), Inches(2.8), Inches(1.2))
    
    tbl_data = [
        ["Source Entity", "Description", "Scale / Coverage", "Update Freq"],
        ["AMFI India", "Scheme AUM growth, total folio counts & Master records", "90 AMC rows, 48 SIP rows", "Monthly"],
        ["mfapi.in JSON API", "Historical daily price data for 40 schemes", "46,000+ price points", "Daily"],
        ["NSE & BSE India", "Benchmark closing index levels (Nifty50, Nifty100, etc.)", "8,050 index price rows", "Daily"],
        ["Simulated Transactions", "Retail transactions (SIP, Lumpsum, Redemption)", "32,778 rows (5,000 investors)", "N/A"]
    ]
    add_table(slide, tbl_data, Inches(0.5), Inches(3.0), Inches(12.33), Inches(3.8), 
              col_widths=[Inches(2.5), Inches(4.83), Inches(3.0), Inches(2.0)])
              
    # -------------------------------------------------------------------------
    # SLIDE 4: SYSTEM ARCHITECTURE
    # -------------------------------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_BG)
    add_header(slide, "5-Layer System Architecture")
    
    # Text box representation of flowchart
    layers = [
        "1. INGESTION LAYER:\nExtract raw CSV file formats from AMFI and request daily live NAVs from JSON API endpoints.",
        "2. PIPELINE ETL:\nParse attributes using Python Pandas. Clean anomalies, forward-fill missing dates, and backfill YoY ratios.",
        "3. SCHEMA STORAGE:\nLoad normalised facts and dimensions into SQLite 3. Build 6 indexes on keys to optimize query execution.",
        "4. ANALYTICS CORE:\nCalculate Sharpe, OLS alpha, and maximum drawdowns. Model cohort sizes and transaction gaps in notebooks.",
        "5. DASHBOARD UI:\nRender results in Streamlit. Implement interactive filter settings, scorecards, and custom recommenders."
    ]
    add_bullet_list(slide, layers, Inches(0.5), Inches(1.5), Inches(12.33), Inches(5.2), font_size=13)
      # -------------------------------------------------------------------------
    # SLIDE 5: EDA HIGHLIGHTS: MARKET GROWTH & FOLIO EXPANSION
    # -------------------------------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_BG)
    add_header(slide, "EDA Highlights: Market Growth & Folio Expansion")
    
    # Insert SIP growth chart
    chart_sip = CHARTS / 'chart_03_sip_inflow_trend.png'
    if chart_sip.exists():
        slide.shapes.add_picture(str(chart_sip), Inches(0.5), Inches(1.3), width=Inches(5.8), height=Inches(4.0))
    else:
        add_textbox(slide, "[SIP Inflow Trend Chart Missing]", Inches(0.5), Inches(2.5), Inches(5.8), Inches(1.0), font_size=14, bold=True)
        
    # Insert Folio growth chart
    chart_folio = CHARTS / 'chart_07_folio_count_growth.png'
    if chart_folio.exists():
        slide.shapes.add_picture(str(chart_folio), Inches(6.8), Inches(1.3), width=Inches(6.0), height=Inches(4.0))
    else:
        add_textbox(slide, "[Folio Growth Chart Missing]", Inches(6.8), Inches(2.5), Inches(6.0), Inches(1.0), font_size=14, bold=True)
        
    insights = [
        "• Monthly SIP inflows grew 2.69x from Jan 2022 (₹11,517 Cr) to Dec 2025 (₹31,002 Cr) showing massive retail participation.",
        "• Total industry folios doubled from 13.26 to 26.12 crore, driven almost entirely by retail equity folio expansion.",
        "• The 26-35 age group represents the largest investor demographic segment, contributing 41% of all transaction volume."
    ]
    add_bullet_list(slide, insights, Inches(0.5), Inches(5.6), Inches(12.33), Inches(1.5), font_size=11)
    
    # -------------------------------------------------------------------------
    # SLIDE 6: FUND PERFORMANCE & SCORECARD
    # -------------------------------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_BG)
    add_header(slide, "Risk-Return Matrix & Top Fund Ranks")
    
    chart_rr = CHARTS / 'chart_08_risk_return_matrix.png'
    if chart_rr.exists():
        slide.shapes.add_picture(str(chart_rr), Inches(0.5), Inches(1.3), width=Inches(6.8), height=Inches(4.2))
    else:
        add_textbox(slide, "[Risk-Return Chart Missing]", Inches(0.5), Inches(2.5), Inches(6.8), Inches(1.0), font_size=14, bold=True)
        
    tbl_sc = [
        ["Rank", "Scheme Name", "Category", "Sharpe", "Score"],
        ["1", "SBI Small Cap Direct", "Small Cap", "0.93", "72.75"],
        ["2", "Kotak Flexicap Regular", "Flexi Cap", "0.85", "71.50"],
        ["3", "Mirae Asset Large Cap Reg", "Large Cap", "1.06", "70.25"],
        ["4", "Quant Active Direct", "Multi Cap", "0.89", "69.10"],
        ["5", "SBI Bluechip Regular", "Large Cap", "0.88", "68.45"]
    ]
    add_textbox(slide, "TOP 5 SCORECARD PICKS (3yr)", Inches(7.6), Inches(1.3), Inches(5.2), Inches(0.4), font_size=14, bold=True, color=RED_ACCENT)
    add_table(slide, tbl_sc, Inches(7.6), Inches(1.8), Inches(5.2), Inches(2.8), col_widths=[Inches(0.6), Inches(2.2), Inches(1.0), Inches(0.7), Inches(0.7)])
    
    rr_insights = [
        "• Small Cap schemes deliver top return rates (20-23% 3yr CAGR) but sit in the high-volatility quadrant.",
        "• Liquid/Debt schemes sit in the bottom-left with low risk (VaR95 as low as -0.02% daily) and low return (5-7%)."
    ]
    add_bullet_list(slide, rr_insights, Inches(7.6), Inches(4.8), Inches(5.2), Inches(2.0), font_size=11)
    
    # -------------------------------------------------------------------------
    # SLIDE 7: ADVANCED RISK & CONCENTRATION
    # -------------------------------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_BG)
    add_header(slide, "Advanced Analytics & Portfolio Tail Risk")
    
    chart_var = CHARTS / 'chart_19_var_distribution.png'
    if chart_var.exists():
        slide.shapes.add_picture(str(chart_var), Inches(0.5), Inches(1.3), width=Inches(5.0), height=Inches(4.5))
    else:
        add_textbox(slide, "[VaR Distribution Missing]", Inches(0.5), Inches(2.5), Inches(5.0), Inches(1.0), font_size=14, bold=True)
        
    chart_hhi = CHARTS / 'chart_22_sector_hhi.png'
    if chart_hhi.exists():
        slide.shapes.add_picture(str(chart_hhi), Inches(5.8), Inches(1.3), width=Inches(4.2), height=Inches(4.5))
    else:
        add_textbox(slide, "[Sector HHI Missing]", Inches(5.8), Inches(2.5), Inches(4.2), Inches(1.0), font_size=14, bold=True)
        
    risk_notes = [
        "Tail Risk & Sector Concentration:",
        "• Small Cap schemes carry 120x more daily tail risk than Liquid schemes (SBI Small Cap Direct daily VaR95 = -2.69%).",
        "• HHI concentration analysis shows all 34 equity portfolios fall below SEBI's 2,500 threshold (Axis Bluechip leads at HHI 2,064)."
    ]
    add_bullet_list(slide, risk_notes, Inches(10.2), Inches(1.8), Inches(2.7), Inches(4.2), font_size=11)
    
    # -------------------------------------------------------------------------
    # SLIDE 8: STREAMLIT ANALYTICS DASHBOARD
    # -------------------------------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_BG)
    add_header(slide, "Streamlit Interactive Analytics Dashboard")
    
    chart_d1 = CHARTS / 'dashboard_page1_industry_overview.png'
    if chart_d1.exists():
        slide.shapes.add_picture(str(chart_d1), Inches(0.5), Inches(1.3), width=Inches(5.8), height=Inches(3.8))
    else:
        add_textbox(slide, "[Dashboard Page 1 Missing]", Inches(0.5), Inches(2.0), Inches(5.8), Inches(1.0), font_size=14, bold=True)
        
    chart_d2 = CHARTS / 'dashboard_page2_fund_performance.png'
    if chart_d2.exists():
        slide.shapes.add_picture(str(chart_d2), Inches(6.8), Inches(1.3), width=Inches(6.0), height=Inches(3.8))
    else:
        add_textbox(slide, "[Dashboard Page 2 Missing]", Inches(6.8), Inches(2.0), Inches(6.0), Inches(1.0), font_size=14, bold=True)
        
    db_caption = [
        "• Real-Time Interactive Web Application built using Streamlit and Plotly.",
        "• 5 Subpages: Industry Overview | Fund Performance | Investor Analytics | SIP Trends | Fund Recommender",
        "• Exposes live sidebar filters: State Selector & Date Range, alongside individual/multi-fund indexed comparisons."
    ]
    add_bullet_list(slide, db_caption, Inches(0.5), Inches(5.4), Inches(12.33), Inches(1.5), font_size=11)
    
    # -------------------------------------------------------------------------
    # SLIDE 9: BEHIND THE SCENES: TECHNICAL CHALLENGES
    # -------------------------------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_BG)
    add_header(slide, "Behind the Scenes: Technical Challenges")
    
    bts_text = [
        "• <b>Resampling NAV Gaps:</b> Resampling daily records to business-day schedule and ffilling gaps threw multiple index errors.",
        "• <b>Sortino Downside Std Dev:</b> Took 3 tries to get the math right. Dividing by negative days instead of total N was a major trap.",
        "• <b>ReportLab Flowables:</b> Margins were a complete nightmare! Balancing tables and charts to prevent overflow took ages.",
        "• <b>Rate Limiting (429):</b> api.mfapi.in kept locking me out during rapid requests. Added retry loops with exponential backoff."
    ]
    add_textbox(slide, "MAJOR HURDLES OVERCOME", Inches(0.5), Inches(1.3), Inches(12.33), Inches(0.4), font_size=14, bold=True, color=RED_ACCENT)
    add_bullet_list(slide, bts_text, Inches(0.5), Inches(1.8), Inches(12.33), Inches(3.0), font_size=12)
    add_textbox(slide, "Note: ReportLab layout is very sensitive. Rebuilt python scripts multiple times to align tables. (Left in as a reminder to refactor later)", Inches(0.5), Inches(5.8), Inches(12.33), Inches(0.8), font_size=10, color=GREY)

    # -------------------------------------------------------------------------
    # SLIDE 10: LESSONS LEARNED & RETROSPECTIVE
    # -------------------------------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_BG)
    add_header(slide, "Lessons Learned & Retrospective")
    
    lessons_bullets = [
        "• <b>Database Normalization (3NF):</b> Denormalizing scheme performance data was a mistake that caused redundancy. Normalizing and using SQL JOINs was a great lesson.",
        "• <b>Clean Ingestion:</b> Separating data retrieval (timeouts/retry limits) from report generation is crucial for stable data pipelines.",
        "• <b>What Didn't Work:</b> Attempted to compile PDF tables dynamically from API data, but it crashed. Separated API fetch and DB loading for stability.",
        "• <b>What I'd Do Differently:</b> If I had more time, I would write proper unit tests for metrics and configure a cloud database (RDS) instead of local SQLite."
    ]
    add_bullet_list(slide, lessons_bullets, Inches(0.5), Inches(1.5), Inches(12.33), Inches(5.0), font_size=12)
    
    # -------------------------------------------------------------------------
    # SLIDE 11: KEY FINDINGS
    # -------------------------------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_BG)
    add_header(slide, "10 Key Project Findings")
    
    key_findings = [
        "1. SIP monthly inflows grew 2.69x, peaking at ₹31,002 crore in December 2025.",
        "2. Retail folios doubled (13.26 to 26.12 crore) driven entirely by equity mutual fund participation.",
        "3. Small Cap schemes deliver top returns (23.39% 3yr CAGR) but carry 120x more daily risk than Liquid funds.",
        "4. SBI Mutual Fund leads AMC assets at ₹12.50 lakh crore (19.8% market share of top-10 AMCs).",
        "5. Liquid mutual funds dominated net category inflows in FY25, drawing ₹4,51,275 crore.",
        "6. The 26-35 age group represents the largest investor demographic cohort (41% transaction count).",
        "7. Average monthly SIP sizes are flat across all age cohorts (ranging ₹10,886 to ₹11,575).",
        "8. T30 cities account for 65.9% of total SIP value, while B30 represents 34.1%.",
        "9. Axis Bluechip has the highest sector concentration (HHI: 2,064) but remains below SEBI's 2,500 threshold.",
        "10. 97.8% of SIP accounts with 6+ historic transactions are at-risk, exhibiting average gaps exceeding 35 days."
    ]
    add_bullet_list(slide, key_findings, Inches(0.5), Inches(1.3), Inches(12.33), Inches(5.6), font_size=11)
    
    # -------------------------------------------------------------------------
    # SLIDE 12: THANK YOU
    # -------------------------------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_BG)
    
    add_textbox(slide, "THANK YOU", Inches(1.0), Inches(2.0), Inches(11.33), Inches(0.6), font_size=32, bold=True, color=RED_ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, "Questions & Discussion", Inches(1.0), Inches(2.8), Inches(11.33), Inches(0.5), font_size=18, color=WHITE, align=PP_ALIGN.CENTER)
    
    add_textbox(slide, "Built with Python · SQLite · Streamlit · Plotly", Inches(1.0), Inches(3.8), Inches(11.33), Inches(0.5), font_size=13, color=GREY, align=PP_ALIGN.CENTER)
    
    add_textbox(slide, "GitHub: github.com/shaikhadibbb/bluestock-mf-capstone\nLocal Dashboard: streamlit run dashboard/app.py\nData Sources: AMFI India · mfapi.in · NSE & BSE India", Inches(1.0), Inches(5.2), Inches(11.33), Inches(1.2), font_size=11, color=WHITE, align=PP_ALIGN.CENTER)
    
    # Save Presentation
    prs.save(str(REPORTS / 'Bluestock_MF_Presentation.pptx'))
    print("Saved Bluestock_MF_Presentation.pptx")
    
if __name__ == '__main__':
    main()
